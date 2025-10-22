from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def set_para_style(p, size=20, bold=False, color=(0,0,0)):
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def add_slide(prs, title, bullets=None, notes=None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    if bullets:
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
        for i, item in enumerate(bullets):
            if isinstance(item, tuple):
                level, text = item
            else:
                level, text = 0, str(item)
            if i == 0:
                p = tf.paragraphs[0]
                p.text = text
            else:
                p = tf.add_paragraph()
                p.text = text
            p.level = level
            set_para_style(p, size=20)

    if notes:
        ns = slide.notes_slide.notes_text_frame
        ns.text = notes
    return slide

def main():
    prs = Presentation()

    # Portada
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Tema 2: jQuery y Google APIs"
    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        "ISW-512 Diseño de Aplicaciones Web – Exposición\n"
        "Carrera: Ingeniería del Software – UTN\n"
        "Profesor: Mtr. Sergio González Salazar\n"
        "Autor: Jafet"
    )

    add_slide(
        prs, "jQuery: definición e historia",
        bullets=[
            "Librería JavaScript (2006) para simplificar DOM, eventos y AJAX.",
            "Popular en la era pre-ES6; aún presente en proyectos y plantillas.",
            "Hoy: JS moderno cubre gran parte; jQuery sigue útil en mantenimiento (legacy)."
        ],
        notes="Enfatiza contexto histórico y por qué todavía lo verán en empresas."
    )

    add_slide(
        prs, "Instalación y sintaxis básica",
        bullets=[
            "Instalación por CDN: https://code.jquery.com/jquery-3.7.1.min.js",
            "Patrón: $(selector).acción()",
            "Encadenamiento: $('.btn').addClass('activa').text('OK')"
        ],
        notes="Muestra en vivo en la consola del navegador una acción simple."
    )

    add_slide(
        prs, "Selectores en jQuery",
        bullets=[
            "#id, .clase, tag",
            "[attr=value] y combinadores (p.ej., 'ul li.item')",
            "Buenas prácticas: selectores específicos y cachear resultados"
        ],
        notes="Conecta con selectores CSS; evitar selectores demasiado genéricos."
    )

    add_slide(
        prs, "Acceso/Manipulación del DOM y eventos",
        bullets=[
            "Lectura/escritura: text(), html(), val(), attr()",
            "Estilos y clases: addClass(), removeClass(), css()",
            "Eventos: on('click', fn), delegación de eventos para elementos dinámicos"
        ],
        notes="Menciona prevención de fugas de memoria removiendo listeners cuando aplique."
    )

    add_slide(
        prs, "Geolocation API del navegador",
        bullets=[
            "navigator.geolocation.getCurrentPosition(success, error, options)",
            "Errores comunes: permiso denegado, no disponible, timeout",
            "Requisitos: HTTPS o localhost"
        ],
        notes="Habla de privacidad y UX al solicitar permisos."
    )

    add_slide(
        prs, "Google Maps JavaScript API",
        bullets=[
            "Carga: .../maps/api/js?key=TU_API_KEY&libraries=geometry",
            "Objetos clave: Map, Marker, InfoWindow",
            "Restricción de API key por HTTP referrers; habilitar billing"
        ],
        notes="Explica cómo proteger la API key y habilitar servicios necesarios."
    )

    add_slide(
        prs, "Cálculo de distancias entre puntos",
        bullets=[
            "google.maps.geometry.spherical.computeDistanceBetween(a, b)",
            "Resultado en metros; conversión sencilla a km",
            "Precisión y formato de salida"
        ],
        notes="Aclara que es distancia en línea recta (no ruta)."
    )

    add_slide(
        prs, "Trazado de rutas (Directions Service)",
        bullets=[
            "DirectionsService + DirectionsRenderer",
            "Parámetros: origin, destination, travelMode (DRIVING/WALKING/etc.)",
            "Manejo de estados y errores del servicio"
        ],
        notes="Demuestra una ruta desde la ubicación actual hacia el negocio."
    )

    add_slide(
        prs, "Google Charts",
        bullets=[
            "Carga: google.charts.load('current', {packages:['corechart']})",
            "DataTable + PieChart/ColumnChart",
            "Uso en dashboards rápidos; alternativa: Chart.js"
        ],
        notes="Presenta un gráfico con la distancia calculada como ejemplo."
    )

    add_slide(
        prs, "Cloud Translation (opcional)",
        bullets=[
            "Servicio de Google Cloud (v2/v3) — requiere proyecto y billing",
            "No exponer claves en el frontend; usar backend intermedio",
            "Casos de uso: internacionalización, soporte, análisis"
        ],
        notes="Menciona brevemente sin demo en vivo para no exponer credenciales."
    )

    add_slide(
        prs, "Demo en vivo",
        bullets=[
            "1) Obtener permiso de geolocalización",
            "2) Marcar posición actual y destino",
            "3) Calcular distancia (geometry)",
            "4) Trazar ruta (Directions)",
            "5) Mostrar gráfico (Google Charts)"
        ],
        notes="Ten preparada una ubicación por defecto en caso de que se niegue el permiso."
    )

    add_slide(
        prs, "Buenas prácticas",
        bullets=[
            "Restringe tu API key y monitoriza uso/costos",
            "Manejo de errores y fallbacks de UX",
            "Performance: limita listeners y actualizaciones del DOM"
        ],
        notes="Incluye métricas básicas y logs para diagnóstico."
    )

    add_slide(
        prs, "Conclusiones",
        bullets=[
            "jQuery agiliza tareas en proyectos existentes",
            "Google Maps/Charts enriquecen la UX con poco código",
            "Siguiente paso: integrar servicios desde backend seguro"
        ]
    )

    add_slide(
        prs, "Recursos",
        bullets=[
            "jQuery API: https://api.jquery.com/",
            "MDN Geolocation: https://developer.mozilla.org/docs/Web/API/Geolocation_API",
            "Maps JS: https://developers.google.com/maps/documentation/javascript",
            "Directions: https://developers.google.com/maps/documentation/javascript/directions",
            "Geometry: https://developers.google.com/maps/documentation/javascript/geometry",
            "Google Charts: https://developers.google.com/chart",
            "Cloud Translation: https://cloud.google.com/translate/docs"
        ],
        notes="Añade enlaces del proyecto/demo si los publicas."
    )

    prs.save("Tema2_jQuery_y_Google_APIs.pptx")
    print("Presentación generada: Tema2_jQuery_y_Google_APIs.pptx")

if __name__ == "__main__":
    main()